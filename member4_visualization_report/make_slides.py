from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x3E, 0x50)
RED = RGBColor(0xE7, 0x4C, 0x3C)
BLUE = RGBColor(0x34, 0x98, 0xDB)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xEC, 0xF0, 0xF1)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=20, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
    return txBox

def add_image(slide, img_path, left, top, width, height=None):
    if os.path.exists(img_path):
        if height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width))
    return

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_textbox(slide, 1.5, 2.0, 10, 1.5, "Traveling Salesman Problem", 44, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.5, 10, 1, "Branch and Bound Algorithm — Performance Analysis", 28, False, RGBColor(0xBD, 0xC3, 0xC7), PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.5, 10, 1, "Group Project Presentation", 22, False, GRAY, PP_ALIGN.CENTER)

# ==================== SLIDE 2: Problem ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "The Traveling Salesman Problem", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Given n cities and distances between each pair, find the shortest route that visits every city exactly once and returns to the start.",
    "NP-hard combinatorial optimization problem.",
    "Real-world applications: logistics, PCB drilling, DNA sequencing, route planning.",
    "This project: Compare Brute Force, Branch & Bound, and Nearest Neighbor algorithms."
], 22, WHITE)

# ==================== SLIDE 3: B&B Theory ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Branch and Bound Algorithm", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Branching: DFS recursive exploration of partial tours.",
    "Bounding: Compute lower bound = current cost + sum of minimum outgoing edges of unvisited cities.",
    "Pruning: If lower bound >= current best cost, discard the entire branch.",
    "Heuristic ordering: Visit nearest neighbors first to find good solutions early.",
    "Result: Significant reduction in search space while guaranteeing optimality."
], 22, WHITE)

# ==================== SLIDE 4: Lower Bound ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Lower Bound & Pruning Logic", 36, True, RED, PP_ALIGN.LEFT)
add_textbox(slide, 0.5, 1.8, 12, 1.5, "LB = current_path_cost + Σ min_outgoing_edge(unvisited_city_i)", 28, True, RGBColor(0xF3, 0x9C, 0x12), PP_ALIGN.CENTER)
add_bullet_slide(slide, 0.5, 3.5, 12, 3.5, [
    "min_outgoing_edge(i) is the shortest edge from city i to any unvisited city (or the start city).",
    "This is an admissible heuristic: it never overestimates the true remaining cost.",
    "If LB >= best_cost → prune (safe: this branch cannot yield a better solution).",
    "More aggressive pruning = smaller search space = faster runtime."
], 20, WHITE)

# ==================== SLIDE 5: Experiment ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Experimental Setup", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Random TSP instances: n = 4, 5, 6, 7, 8, 9, 10 cities.",
    "5 random instances per size (coordinates uniformly distributed in 1000x1000 grid).",
    "Three algorithms tested: Brute Force (control), Branch & Bound, Nearest Neighbor.",
    "Metrics: runtime (s), visited nodes, pruned branches, pruning rate, solution quality.",
    "All tests run on identical datasets for fair comparison."
], 22, WHITE)

# ==================== SLIDE 6: Results Runtime ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_textbox(slide, 0.5, 0.3, 12, 1, "Results: Runtime vs Problem Size", 36, True, DARK, PP_ALIGN.LEFT)
add_image(slide, r'D:\0work\www\成员4\runtime_by_cities.png', 1, 1.5, 11, 5.5)

# ==================== SLIDE 7: Search Space ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_textbox(slide, 0.5, 0.3, 12, 1, "Results: Search Space Reduction", 36, True, DARK, PP_ALIGN.LEFT)
add_image(slide, r'D:\0work\www\成员4\nodes_by_cities.png', 0.5, 1.5, 6, 4)
add_image(slide, r'D:\0work\www\成员4\nodes_bar_comparison.png', 6.8, 1.5, 6, 4)

# ==================== SLIDE 8: Pruning ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_textbox(slide, 0.5, 0.3, 12, 1, "Results: Pruning Efficiency", 36, True, DARK, PP_ALIGN.LEFT)
add_image(slide, r'D:\0work\www\成员4\pruning_rate_by_cities.png', 0.5, 1.5, 6, 4)
add_image(slide, r'D:\0work\www\成员4\pruned_branches_by_cities.png', 6.8, 1.5, 6, 4)

# ==================== SLIDE 9: Quality ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_textbox(slide, 0.5, 0.3, 12, 1, "Results: Route Quality (NN vs Optimal)", 36, True, DARK, PP_ALIGN.LEFT)
add_image(slide, r'D:\0work\www\成员4\quality_comparison.png', 1.5, 1.5, 10, 5)

# ==================== SLIDE 10: Summary - Algorithms Recap ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Summary: Three Algorithms Compared", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Brute Force — Exact baseline. Enumerates all (n-1)!/2 permutations. Guarantees optimality but O(n!). Impractical beyond n=10.",
    "Branch & Bound — Exact with pruning. DFS + lower bound evaluation. Cuts branches when LB >= best_cost. Optimal, far more efficient.",
    "Nearest Neighbor — Greedy heuristic. Always picks the closest unvisited city. O(n²) speed but 2–10% above optimal on average.",
    "Takeaway: B&B strikes the best balance between correctness and efficiency for exact TSP solving at small to medium scale."
], 22, WHITE)

# ==================== SLIDE 11: Summary - Key Results ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Summary: Key Results", 36, True, RED, PP_ALIGN.LEFT)
# Results table
table_data = [
    ["N", "B&B Runtime", "BF Runtime", "B&B Nodes", "BF Nodes", "Pruning Rate"],
    ["4", "0.02ms", "0.006ms", "15", "6", "3.8%"],
    ["7", "0.66ms", "0.28ms", "428", "720", "78.1%"],
    ["9", "6.5ms", "14.3ms", "4,107", "40,320", "96.3%"],
    ["10", "23ms", "N/A", "12,472", "1,814,400", "98.7%"],
]
rows, cols = len(table_data), len(table_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.5)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xDF, 0xE6, 0xE9)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

add_textbox(slide, 0.8, 5.5, 11, 1.5,
    "Key insight: B&B prunes 98.7% of search space at n=10. It becomes faster than BF from n=7 onward. "
    "For n=10, BF is infeasible (1.8M paths) while B&B completes in 23ms.",
    18, False, GRAY, PP_ALIGN.LEFT)

# ==================== SLIDE 12: Summary - Conclusions ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Summary: Conclusions", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Branch and Bound is the best choice among these three methods when we need an exact TSP solution but also want better efficiency than simple exhaustive search.",
    "The pruning mechanism becomes more effective as problem size increases — the lower bound provides increasing value for larger instances.",
    "Nearest Neighbor is the fastest but may produce longer routes. Suitable when speed is critical and optimality can be sacrificed.",
    "Brute Force is accurate but does not scale. Its exponential growth makes it unusable beyond 10–12 cities.",
    "Limitation: B&B remains exponential; for n > 20 it becomes impractical. Tighter bounds (MST, Held-Karp) could improve pruning."
], 20, WHITE)

# ==================== SLIDE 13: Key Findings ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Key Findings", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "B&B finds provably optimal solutions (same as brute force) for all tested instances.",
    "At n=10, B&B prunes 98.7% of the search space — visiting only 1.3% of all possible nodes.",
    "B&B is ~2.2× faster than brute force at n=9; gap widens with larger n.",
    "Nearest Neighbor runs in near-zero time but solutions are 2-10% above optimal.",
    "B&B is practical for n ≤ 14; beyond that, heuristics are needed."
], 22, WHITE)

# ==================== SLIDE 14: Conclusion ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.5, 0.3, 12, 1, "Conclusion & Future Work", 36, True, RED, PP_ALIGN.LEFT)
add_bullet_slide(slide, 0.5, 1.5, 12, 5.5, [
    "Branch and Bound is an effective exact algorithm for small TSP instances.",
    "The minimum-edge-sum lower bound provides substantial pruning even at modest n.",
    "Future work: Compare with MST-based and Held-Karp lower bounds.",
    "Future work: Implement parallel B&B for larger instances.",
    "Future work: Hybrid approaches combining B&B with metaheuristics."
], 22, WHITE)

# ==================== SLIDE 15: Thanks ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 1.5, 2.5, 10, 1.5, "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.5, 10, 1, "Questions?", 28, False, GRAY, PP_ALIGN.CENTER)

output_path = r'D:\0work\www\成员4\tsp_bb_presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
